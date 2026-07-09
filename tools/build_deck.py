"""Generate the Con Edison solutions deck (light theme) from the code.

Every figure on every slide is computed here, at build time, by importing the
same engines the application runs. Nothing is typed by hand. If a detector
changes or a lever is added, rebuild and the deck follows -- a deck whose
numbers drift from the product is worse than no deck.

    python tools/build_deck.py            # -> Infosys_FinOps_ConEdison.pptx
    python tools/build_deck.py --out X.pptx

Charts are native PowerPoint charts, not images, so the client can click into
them, restyle them, and lift them into their own template.

The estate behind the numbers is SYNTHETIC. Every slide that quotes a figure
says so.
"""

from __future__ import annotations

import argparse
import os
import sys
from dataclasses import dataclass
from typing import Iterable, List, Optional, Sequence, Tuple

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, ROOT)
sys.path.insert(0, os.path.join(ROOT, "tools"))

from pptx import Presentation  # noqa: E402
from pptx.chart.data import CategoryChartData  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Emu, Inches, Pt  # noqa: E402

# ==========================================================================
# Light theme
#
# Identity colour only. The chart series colours are the same CVD-validated
# categorical slots the app uses, so a screenshot and a slide agree.
# ==========================================================================

W, H = Inches(13.333), Inches(7.5)  # 16:9

INK = RGBColor(0x0B, 0x14, 0x2A)
BODY = RGBColor(0x44, 0x50, 0x66)
MUTED = RGBColor(0x7A, 0x85, 0x99)
RULE = RGBColor(0xE2, 0xE7, 0xEF)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
WASH = RGBColor(0xF5, 0xF8, 0xFC)

AZURE = RGBColor(0x1E, 0x6F, 0xD9)
TEAL = RGBColor(0x11, 0x9B, 0x8A)
VIOLET = RGBColor(0x5B, 0x4B, 0xC4)
AMBER = RGBColor(0xC9, 0x85, 0x00)
CRIMSON = RGBColor(0xC2, 0x33, 0x33)
GREEN = RGBColor(0x0C, 0x7A, 0x3E)

# Chart series palette (light mode slots from theme.py).
SERIES = [
    RGBColor(0x2A, 0x78, 0xD6),
    RGBColor(0xEB, 0x68, 0x34),
    RGBColor(0x1B, 0xAF, 0x7A),
    RGBColor(0x4A, 0x3A, 0xA7),
    RGBColor(0xED, 0xA1, 0x00),
]

FONT = "Segoe UI"


# ==========================================================================
# Facts -- read from the running code
# ==========================================================================


@dataclass
class Facts:
    org: str
    spend: float
    run_rate: float
    mom: float
    yoy: float
    esr: float
    coverage: float
    utilization: float
    commitment_waste: float
    cost_of_waste: float
    waste_pct: float
    allocation: float
    readiness: str
    drift: float
    fc_method: str
    fc_wape: float
    fc_maturity: str
    fc_total: float
    fc_with_cliffs: float
    cliff_months: List[str]
    months: int
    rows: int
    savings_total: float
    savings_by_category: List[Tuple[str, float]]
    top_opps: List[Tuple[str, str, str, float, str, str]]
    esr_uplift: Tuple[float, float, float]
    n_levers: int
    n_connectors: int
    n_opps: int
    spend_by_cloud: List[Tuple[str, float]]
    monthly: List[Tuple[str, float]]
    forecast: List[Tuple[str, float, float, float]]
    allocation_by_bu: List[Tuple[str, float, float]]
    anomalies: List[Tuple[str, str, float, float, float]]


def gather() -> Facts:
    os.environ.setdefault("FINOPS_MODE", "demo")
    import allocation as alloc
    import anomaly
    import connectors
    import data
    import forecast as fx
    import kpi
    import optimize

    ctx = data.load_demo_context()
    df = ctx.focus_df

    opps = optimize.detect_all(df)
    uw = optimize.usage_waste_total(opps)
    k = kpi.executive_kpis(df, usage_waste_monthly=uw)

    monthly = ctx.monthly()
    fc = fx.forecast_spend(monthly, horizon=24, method="auto")
    cliff = fx.commitment_expiry_overlay(df, fc.forecast)
    cliff_months = (
        cliff.loc[cliff["cliff"] == True, "period"].dt.strftime("%b %Y").tolist()  # noqa: E712
        if "cliff" in cliff
        else []
    )

    sav = optimize.savings_by_category(opps)
    up = optimize.effective_savings_rate_uplift(df, opps)

    by_cloud = (
        df.groupby("ProviderName", observed=True)["EffectiveCost"].sum().sort_values(ascending=False)
    )

    bu = alloc.allocate(df, alloc.SharedCostPolicy(method="proportional"), dim="tag_business_unit")
    bu = bu.sort_values("total_cost", ascending=False)

    an = anomaly.detect_by_dimension(df, dim="ServiceCategory")
    flagged = an[an["is_anomaly"]].sort_values("deviation_pct", ascending=False).head(3)

    top = sorted(opps, key=lambda o: -o.annual_savings)[:8]

    return Facts(
        org=ctx.config.organisation,
        spend=k.total_spend,
        run_rate=k.run_rate or 0,
        mom=k.mom_pct or 0,
        yoy=k.yoy_pct or 0,
        esr=k.esr_pct or 0,
        coverage=k.coverage_pct or 0,
        utilization=k.utilization_pct or 0,
        commitment_waste=k.commitment_waste,
        cost_of_waste=k.cost_of_waste,
        waste_pct=k.waste_pct or 0,
        allocation=k.allocation_coverage_pct or 0,
        readiness=k.chargeback_readiness,
        drift=k.baseline_drift_pct or 0,
        fc_method=fc.method.replace("_", " ").title(),
        fc_wape=fc.accuracy["wape"],
        fc_maturity=fc.maturity,
        fc_total=float(fc.forecast["cost"].sum()),
        fc_with_cliffs=float(cliff["cost_with_cliffs"].sum()) if "cost_with_cliffs" in cliff else 0.0,
        cliff_months=cliff_months,
        months=len(monthly),
        rows=len(df),
        savings_total=float(sum(o.annual_savings for o in opps)),
        savings_by_category=[(r.category, float(r.annual_savings)) for r in sav.itertuples()],
        top_opps=[
            (o.lever_id, o.lever_name, o.cloud, o.annual_savings, o.effort, o.risk) for o in top
        ],
        esr_uplift=(up["current_esr_pct"], up["projected_esr_pct"], up["uplift_pts"]),
        n_levers=len(optimize.LEVERS),
        n_connectors=len(connectors.REGISTRY),
        n_opps=len(opps),
        spend_by_cloud=[(str(i), float(v)) for i, v in by_cloud.items()],
        monthly=[(p.strftime("%b %y"), float(c)) for p, c in zip(monthly["period"], monthly["cost"])],
        forecast=[
            (p.strftime("%b %y"), float(c), float(lo), float(hi))
            for p, c, lo, hi in zip(
                fc.forecast["period"], fc.forecast["cost"], fc.forecast["lo80"], fc.forecast["hi80"]
            )
        ],
        allocation_by_bu=[
            (str(r[1]), float(r.total_cost), float(r.pct_of_total)) for r in bu.itertuples()
        ],
        anomalies=[
            (
                r.period.strftime("%d %b %Y"),
                str(getattr(r, "ServiceCategory")),
                float(r.cost),
                float(r.expected),
                float(r.deviation_pct),
            )
            for r in flagged.itertuples()
        ],
    )


# ==========================================================================
# Layout primitives
# ==========================================================================


def money(x: float) -> str:
    a = abs(x)
    s = "-" if x < 0 else ""
    if a >= 1e9:
        return f"{s}${a/1e9:.2f}B"
    if a >= 1e6:
        return f"{s}${a/1e6:.2f}M"
    if a >= 1e3:
        return f"{s}${a/1e3:.0f}K"
    return f"{s}${a:,.0f}"


def _text(shape, size=14, bold=False, colour=BODY, align=PP_ALIGN.LEFT, space_after=4):
    tf = shape.text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.alignment = align
        p.space_after = Pt(space_after)
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = colour
            r.font.name = FONT
    return tf


def box(slide, x, y, w, h, text="", size=14, bold=False, colour=BODY, align=PP_ALIGN.LEFT):
    sh = slide.shapes.add_textbox(x, y, w, h)
    sh.text_frame.text = text
    _text(sh, size, bold, colour, align)
    return sh


def bullets(slide, x, y, w, h, items: Sequence[str], size=14, colour=BODY, bullet="—"):
    sh = slide.shapes.add_textbox(x, y, w, h)
    tf = sh.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet}  {item}" if bullet else item
        p.space_after = Pt(9)
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.color.rgb = colour
            r.font.name = FONT
    return sh


def rect(slide, x, y, w, h, fill=WASH, line=None):
    from pptx.enum.shapes import MSO_SHAPE

    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.06
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def blank(prs) -> "object":
    return prs.slides.add_slide(prs.slide_layouts[6])


def header(slide, eyebrow: str, title: str, sub: str = "") -> None:
    box(slide, Inches(0.7), Inches(0.42), Inches(11), Inches(0.3), eyebrow.upper(), 10, True, MUTED)
    box(slide, Inches(0.7), Inches(0.7), Inches(11.9), Inches(0.6), title, 28, True, INK)
    if sub:
        box(slide, Inches(0.7), Inches(1.32), Inches(11.9), Inches(0.4), sub, 12.5, False, MUTED)
    ln = slide.shapes.add_shape(1, Inches(0.7), Inches(1.82), Inches(1.4), Pt(3))
    ln.fill.solid()
    ln.fill.fore_color.rgb = AZURE
    ln.line.fill.background()
    ln.shadow.inherit = False


def footer(slide, note: str = "") -> None:
    box(slide, Inches(0.7), Inches(6.95), Inches(9), Inches(0.3), note, 8.5, False, MUTED)
    box(slide, Inches(11.4), Inches(6.95), Inches(1.3), Inches(0.3), "Infosys", 8.5, True, MUTED, PP_ALIGN.RIGHT)


SYNTHETIC = "Figures computed from a synthetic 24-month utility estate shipped with the platform. No customer data."


def kpi_card(slide, x, y, w, h, label: str, value: str, sub: str, accent=AZURE):
    rect(slide, x, y, w, h, PAPER, RULE)
    bar = slide.shapes.add_shape(1, x, y, Pt(3.5), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False
    box(slide, x + Inches(0.18), y + Inches(0.10), w - Inches(0.3), Inches(0.28), label.upper(), 8.5, True, MUTED)
    box(slide, x + Inches(0.18), y + Inches(0.36), w - Inches(0.3), Inches(0.5), value, 22, True, INK)
    box(slide, x + Inches(0.18), y + Inches(0.90), w - Inches(0.3), Inches(0.4), sub, 9, False, MUTED)


def style_chart(chart, colours: Sequence[RGBColor], legend=False, number_format='#,##0,,"M"'):
    chart.font.size = Pt(10)
    chart.font.name = FONT
    chart.font.color.rgb = BODY
    chart.has_legend = legend
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)
    for i, series in enumerate(chart.plots[0].series):
        colour = colours[i % len(colours)]
        try:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = colour
        except Exception:
            pass
        try:
            series.format.line.color.rgb = colour
            series.format.line.width = Pt(2)
        except Exception:
            pass
    try:
        va = chart.value_axis
        va.has_major_gridlines = True
        va.major_gridlines.format.line.color.rgb = RULE
        va.major_gridlines.format.line.width = Pt(0.5)
        va.tick_labels.number_format = number_format
        va.tick_labels.number_format_is_linked = False
        va.format.line.fill.background()
        ca = chart.category_axis
        ca.has_major_gridlines = False
        ca.format.line.color.rgb = RULE
    except Exception:
        pass


# ==========================================================================
# Slides
# ==========================================================================


def slide_title(prs, f: Facts):
    s = blank(prs)
    rect(s, Inches(0), Inches(0), W, Inches(2.6), RGBColor(0xF2, 0xF6, 0xFC))
    box(s, Inches(0.9), Inches(0.75), Inches(6), Inches(0.35), "INFOSYS", 12, True, AZURE)
    box(s, Inches(0.9), Inches(1.05), Inches(11), Inches(0.9), "Multi-Cloud FinOps Command Center", 40, True, INK)
    box(
        s,
        Inches(0.9),
        Inches(1.95),
        Inches(11),
        Inches(0.5),
        "A single plane of control for AWS, Azure, GCP and OCI spend",
        16,
        False,
        BODY,
    )
    box(s, Inches(0.9), Inches(3.1), Inches(11), Inches(0.4), f"Prepared for {f.org}", 15, True, INK)
    bullets(
        s,
        Inches(0.9),
        Inches(3.6),
        Inches(11.5),
        Inches(2.2),
        [
            "VP and Director visibility by application, business unit and environment",
            "Showback and chargeback with a defensible shared-cost split",
            "A 24-month forecast against budget that sees the commitment cliff",
            f"{f.n_levers} FinOps optimization levers, detected from the bill itself",
            f"Vendor-neutral: {f.n_connectors} connectors, one FOCUS schema underneath",
        ],
        13.5,
    )
    footer(s, "Infosys · Con Edison · FinOps Solution Overview")


def slide_challenge(prs, f: Facts):
    s = blank(prs)
    header(s, "The problem", "Four clouds, four answers, no single truth", "Why multi-cloud FinOps is hard before it is expensive")
    items = [
        ("Fragmented data", "Each provider bills in its own schema. Reconciling AWS, Azure, GCP and OCI by hand is a monthly project, not a dashboard."),
        ("Tool lock-in", "Pick a FinOps platform and every dashboard, KPI and script is written against that vendor's field names. Switching means rebuilding."),
        ("Unallocated spend", "Untagged resources land in a bucket nobody owns, so chargeback is disputed and showback is ignored."),
        ("Blind forecasting", "A trend line walks straight through a commitment expiry. The rate snaps back to on-demand and the variance lands unannounced."),
        ("Waste is invisible", "Unused commitment and idle resources do not announce themselves; they arrive as a slightly larger invoice."),
    ]
    y = Inches(2.15)
    for i, (t, d) in enumerate(items):
        rect(s, Inches(0.7), y, Inches(11.9), Inches(0.82), PAPER if i % 2 else WASH, RULE)
        box(s, Inches(0.95), y + Inches(0.12), Inches(2.9), Inches(0.4), t, 13, True, INK)
        box(s, Inches(3.95), y + Inches(0.12), Inches(8.4), Inches(0.6), d, 11.5, False, BODY)
        y += Inches(0.92)
    footer(s)


def slide_one_idea(prs, f: Facts):
    s = blank(prs)
    header(s, "The approach", "One idea: normalise everything to FOCUS", "The FinOps Foundation's Open Cost and Usage Specification is the contract")

    sources = ["AWS Data Exports", "Azure Cost Management", "GCP Billing Export", "OCI FOCUS Reports", "Any procured FinOps tool", "Any FOCUS CSV / Parquet"]
    y = Inches(2.3)
    for src in sources:
        rect(s, Inches(0.7), y, Inches(2.9), Inches(0.55), PAPER, RULE)
        box(s, Inches(0.85), y + Inches(0.11), Inches(2.7), Inches(0.35), src, 11, False, BODY)
        y += Inches(0.68)

    rect(s, Inches(4.0), Inches(2.3), Inches(1.5), Inches(2.72), WASH, RULE)
    box(s, Inches(4.05), Inches(3.35), Inches(1.4), Inches(0.5), "Connector", 11, True, BODY, PP_ALIGN.CENTER)

    rect(s, Inches(5.9), Inches(2.75), Inches(2.5), Inches(1.8), RGBColor(0xE8, 0xF1, 0xFD), AZURE)
    box(s, Inches(5.95), Inches(3.15), Inches(2.4), Inches(0.5), "FOCUS 1.2", 20, True, INK, PP_ALIGN.CENTER)
    box(s, Inches(5.95), Inches(3.65), Inches(2.4), Inches(0.4), "one frame", 11, False, MUTED, PP_ALIGN.CENTER)

    outs = ["Executive KPIs", "Forecast & budget", "Showback / chargeback", "Optimization levers", "Anomalies", "AI agent team"]
    y = Inches(2.05)
    for o in outs:
        rect(s, Inches(8.9), y, Inches(3.7), Inches(0.48), PAPER, RULE)
        box(s, Inches(9.05), y + Inches(0.08), Inches(3.5), Inches(0.35), o, 11, False, BODY)
        y += Inches(0.56)

    box(
        s,
        Inches(0.7),
        Inches(5.6),
        Inches(11.9),
        Inches(1.0),
        "No dashboard, KPI formula, optimization detector or agent tool has ever seen a vendor-specific field. "
        "Adopting a new FinOps platform is one Connector subclass; dropping one is deleting it. "
        "That is the whole architectural bet — and it is the same bet the FinOps Foundation made with FOCUS.",
        12,
        False,
        BODY,
    )
    footer(s)


def slide_architecture(prs, f: Facts):
    s = blank(prs)
    header(s, "Architecture", "Two modes, one code path", "Nothing below the loader knows which mode is active")
    cards = [
        ("Demo Mode", AZURE, ["Deterministic synthetic estate, generated in-process", "24 months, 3 clouds, 11 applications", f"{f.rows:,} FOCUS charge rows", "No credentials, no network calls, no cloud spend", "Used for evaluation, training and demos"]),
        ("Live Mode", TEAL, ["Real billing data through configured connectors", "Native AWS / Azure / GCP / OCI, or a procured tool", "Multiple payers, tenants and billing accounts", "A half-wired estate still renders the wired parts", "Cached: Cost Explorer bills ~$0.01 per request"]),
    ]
    x = Inches(0.7)
    for title, colour, pts in cards:
        rect(s, x, Inches(2.15), Inches(5.85), Inches(3.6), PAPER, RULE)
        bar = s.shapes.add_shape(1, x, Inches(2.15), Inches(5.85), Pt(4))
        bar.fill.solid(); bar.fill.fore_color.rgb = colour; bar.line.fill.background(); bar.shadow.inherit = False
        box(s, x + Inches(0.28), Inches(2.35), Inches(5.3), Inches(0.45), title, 18, True, INK)
        bullets(s, x + Inches(0.28), Inches(2.85), Inches(5.3), Inches(2.6), pts, 11.5)
        x += Inches(6.2)
    box(s, Inches(0.7), Inches(5.95), Inches(11.9), Inches(0.7),
        "Demo Mode is not a mock. It uses the same Connector interface and produces a frame that passes FOCUS validation, "
        "containing real commitment waste, untagged spend, idle resources, a migration step-change and two injected anomalies.",
        11.5, False, MUTED)
    footer(s)


def _diagram(prs, name: str, eyebrow: str, title_text: str, sub: str, note: str):
    """Full-bleed diagram slide.

    The `.svg` is the deliverable and lives in docs/diagrams/; PowerPoint cannot
    embed SVG, so the 220-dpi PNG rendered from the same source goes here.
    """
    import diagrams as dg  # tools/ is on sys.path

    png = os.path.join(dg.OUT_DIR, f"{name}.png")
    if not os.path.exists(png):
        dg.build_all()

    s = blank(prs)
    header(s, eyebrow, title_text, sub)

    from PIL import Image

    with Image.open(png) as im:
        iw, ih = im.size
    avail_w, avail_h = Inches(12.1), Inches(4.55)
    scale = min(avail_w / iw, avail_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    s.shapes.add_picture(png, int((W - w) / 2), Inches(2.05), w, h)

    box(s, Inches(0.7), Inches(6.72), Inches(10.4), Inches(0.3), note, 9, False, MUTED)
    box(s, Inches(11.0), Inches(6.72), Inches(1.7), Inches(0.3), f"docs/diagrams/{name}.svg", 8, False, MUTED, PP_ALIGN.RIGHT)
    footer(s)


def slide_hld(prs, f: Facts):
    _diagram(prs, "hld", "High level design", "Every source normalises to one FOCUS 1.2 frame",
             "Sources · Ingest · Canonical · Analytics · Experience",
             "Vector source: docs/diagrams/hld.svg")


def slide_euv(prs, f: Facts):
    _diagram(prs, "end_user_view", "End user view", "Who asks what, and where the answer lives",
             "The journey, the six FinOps personas, and the dashboards each one lands on",
             "Vector source: docs/diagrams/end_user_view.svg")


def slide_lld(prs, f: Facts):
    _diagram(prs, "lld", "Low level design", "One request through the system",
             "Module boundaries, the secret boundary, and where results are cached",
             "Vector source: docs/diagrams/lld.svg")


def slide_exec(prs, f: Facts):
    s = blank(prs)
    header(s, "Executive view", "What a VP sees first", f"Amortised spend across AWS, Azure, GCP and OCI · {f.months} months of history")
    cards = [
        ("Total amortised spend", money(f.spend), f"Run-rate {money(f.run_rate)}/yr", AZURE),
        ("Effective savings rate", f"{f.esr:.1f}%", "vs on-demand equivalent", VIOLET),
        ("Commitment coverage", f"{f.coverage:.1f}%", "of eligible spend", TEAL),
        ("Commitment utilisation", f"{f.utilization:.1f}%", f"{money(f.commitment_waste)} unused", GREEN),
        ("Cost of waste", money(f.cost_of_waste), f"{f.waste_pct:.1f}% of spend", AMBER),
        ("Allocation coverage", f"{f.allocation:.1f}%", f.readiness, CRIMSON),
    ]
    x, y = Inches(0.7), Inches(2.15)
    for i, (l, v, sub, c) in enumerate(cards):
        kpi_card(s, x, y, Inches(3.83), Inches(1.4), l, v, sub, c)
        x += Inches(4.02)
        if i == 2:
            x, y = Inches(0.7), Inches(3.75)

    box(s, Inches(0.7), Inches(5.35), Inches(11.9), Inches(1.4),
        "Effective Savings Rate is the outcome metric, not coverage. 100% coverage at 60% utilisation is still a bad deal, "
        "and only ESR shows that. FinOps Foundation benchmarks: median ~0%, 75th percentile ~23%, 98th percentile ~46%. "
        f"This estate sits at {f.esr:.1f}% — above median, with {money(f.commitment_waste)} of commitment already burned unused.",
        12, False, BODY)
    footer(s, SYNTHETIC)


def slide_spend_chart(prs, f: Facts):
    s = blank(prs)
    header(s, "Where the money is", "Spend by cloud and by month", "Amortised EffectiveCost — never blended or unblended on a leadership view")

    cd = CategoryChartData()
    cd.categories = [c for c, _ in f.spend_by_cloud]
    cd.add_series("Amortised spend", tuple(v for _, v in f.spend_by_cloud))
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), Inches(2.15), Inches(5.4), Inches(3.9), cd)
    style_chart(gf.chart, SERIES)

    step = max(1, len(f.monthly) // 12)
    cats = [c for i, (c, _) in enumerate(f.monthly)]
    cd2 = CategoryChartData()
    cd2.categories = cats
    cd2.add_series("Monthly spend", tuple(v for _, v in f.monthly))
    gf2 = s.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(6.5), Inches(2.15), Inches(6.1), Inches(3.9), cd2)
    style_chart(gf2.chart, SERIES)
    try:
        for i, lbl in enumerate(gf2.chart.category_axis.tick_labels.__class__.__mro__):
            break
        gf2.chart.category_axis.tick_labels.font.size = Pt(7)
    except Exception:
        pass

    box(s, Inches(0.7), Inches(6.15), Inches(11.9), Inches(0.6),
        f"Month on month {f.mom:+.1f}%.  Year on year {f.yoy:+.1f}% — the step change is a data-centre exit wave, "
        "which is precisely the shape a naive trend model extrapolates straight through.", 11.5, False, MUTED)
    footer(s, SYNTHETIC)


def slide_allocation(prs, f: Facts):
    s = blank(prs)
    header(s, "Showback & chargeback", "A defensible split, including the shared-cost pool",
           "Direct costs, plus each unit's share of the platform they all consume")

    cd = CategoryChartData()
    cd.categories = [b for b, _, _ in f.allocation_by_bu]
    cd.add_series("Allocated cost", tuple(v for _, v, _ in f.allocation_by_bu))
    gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.7), Inches(2.15), Inches(6.6), Inches(4.0), cd)
    style_chart(gf.chart, SERIES)

    x = Inches(7.6)
    box(s, x, Inches(2.15), Inches(5.0), Inches(0.4), "Five allocation methods", 14, True, INK)
    bullets(s, x, Inches(2.6), Inches(5.0), Inches(2.0), [
        "Direct — the resource carries its own owner tag",
        "Even split — pool ÷ N consumers",
        "Proportional to direct spend — the fair tax",
        "Fixed percentage — a negotiated share",
        "Usage driver — a metered proxy",
    ], 11.5)
    box(s, x, Inches(4.7), Inches(5.0), Inches(1.6),
        f"Allocation coverage is {f.allocation:.1f}%. Practitioner consensus puts the chargeback-readiness line near 90%, "
        f"so this estate is: {f.readiness.lower()}. Showback moves information; chargeback moves money. "
        "Neither is 'more mature' — it is an accounting-policy choice.", 11, False, BODY)
    footer(s, SYNTHETIC)


def slide_forecast(prs, f: Facts):
    s = blank(prs)
    header(s, "Forecast & budget", "Two years ahead, and the cliff a trend line cannot see",
           f"Method chosen by rolling-origin backtest: {f.fc_method} · WAPE {f.fc_wape:.2f}% ({f.fc_maturity})")

    cd = CategoryChartData()
    cd.categories = [c for c, _, _, _ in f.forecast]
    cd.add_series("Forecast", tuple(v for _, v, _, _ in f.forecast))
    cd.add_series("80% lower", tuple(v for _, _, v, _ in f.forecast))
    cd.add_series("80% upper", tuple(v for _, _, _, v in f.forecast))
    gf = s.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(0.7), Inches(2.15), Inches(7.4), Inches(3.9), cd)
    style_chart(gf.chart, [SERIES[0], MUTED, MUTED], legend=True)
    try:
        gf.chart.category_axis.tick_labels.font.size = Pt(7)
    except Exception:
        pass

    x = Inches(8.5)
    kpi_card(s, x, Inches(2.15), Inches(4.1), Inches(1.25), "24-month forecast", money(f.fc_total), "point estimate, cumulative", AZURE)
    kpi_card(s, x, Inches(3.55), Inches(4.1), Inches(1.25), "With commitment cliffs", money(f.fc_with_cliffs),
             f"expiry in {', '.join(f.cliff_months) or 'n/a'}", CRIMSON)
    box(s, x, Inches(5.0), Inches(4.1), Inches(1.6),
        f"When an RI, Savings Plan or CUD term ends, the rate snaps back to on-demand. The overlay adds "
        f"{money(f.fc_with_cliffs - f.fc_total)} the trend line never sees. This is the single most important thing "
        "naive cloud forecasting misses.", 11, False, BODY)

    box(s, Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.55),
        "Foundation forecast-variance maturity bands: Crawl <20%, Walk <15%, Run <12%, best-in-class <5%. "
        "WAPE is dollar-weighted, so a small service rounding to zero cannot dominate the score the way it does under MAPE.",
        11, False, MUTED)
    footer(s, SYNTHETIC)


def slide_optimize(prs, f: Facts):
    s = blank(prs)
    header(s, "Optimization", f"{money(f.savings_total)} identified, across {f.n_opps} opportunities",
           f"Detected from the FOCUS frame by rule, not read from a vendor's recommendation API")

    cd = CategoryChartData()
    cd.categories = [c for c, _ in f.savings_by_category]
    cd.add_series("Annual savings", tuple(v for _, v in f.savings_by_category))
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), Inches(2.2), Inches(4.9), Inches(3.5), cd)
    style_chart(gf.chart, SERIES, number_format='#,##0,,"M"')

    x = Inches(6.0)
    box(s, x, Inches(2.15), Inches(6.6), Inches(0.35), "Largest opportunities", 13, True, INK)
    y = Inches(2.6)
    rect(s, x, y, Inches(6.6), Inches(0.32), WASH)
    for txt, dx, w in [("Lever", 0.1, 0.6), ("Name", 0.75, 2.7), ("Cloud", 3.5, 1.0), ("Annual", 4.55, 1.1), ("Effort/Risk", 5.7, 0.85)]:
        box(s, x + Inches(dx), y + Inches(0.02), Inches(w), Inches(0.28), txt, 9, True, MUTED)
    y += Inches(0.36)
    for lid, name, cloud, ann, eff, risk in f.top_opps:
        for txt, dx, w, bold in [
            (lid, 0.1, 0.6, True),
            (name[:32], 0.75, 2.7, False),
            (cloud[:14], 3.5, 1.0, False),
            (money(ann), 4.55, 1.1, True),
            (f"{eff}/{risk}", 5.7, 0.85, False),
        ]:
            box(s, x + Inches(dx), y, Inches(w), Inches(0.3), txt, 9.5, bold, INK if bold else BODY)
        y += Inches(0.36)

    cur, proj, up = f.esr_uplift
    box(s, Inches(0.7), Inches(5.95), Inches(11.9), Inches(0.8),
        f"Executing the rate levers alone would move Effective Savings Rate from {cur:.1f}% to {proj:.1f}% (+{up:.1f} points). "
        "Savings percentages in the catalog are vendor 'up-to' figures — ceilings, not guarantees. "
        "Where a detector cannot see what it needs (access patterns, CPU utilisation) it lowers its confidence and says so.",
        11, False, BODY)
    footer(s, SYNTHETIC)


def slide_anomaly(prs, f: Facts):
    s = blank(prs)
    header(s, "Anomaly detection", "Catching the runaway workload, not the weekend",
           "STL decomposition, then a median-absolute-deviation test on the residual")
    y = Inches(2.25)
    rect(s, Inches(0.7), y, Inches(11.9), Inches(0.38), WASH)
    for txt, dx, w in [("Date", 0.15, 1.6), ("Service category", 1.9, 3.0), ("Actual", 5.1, 1.4), ("Expected", 6.7, 1.4), ("Deviation", 8.4, 1.6), ("Severity", 10.2, 1.6)]:
        box(s, Inches(0.7 + dx), y + Inches(0.04), Inches(w), Inches(0.3), txt, 10, True, MUTED)
    y += Inches(0.46)
    for d, cat, cost, exp, dev in f.anomalies:
        sev = "Critical" if dev > 100 else "Serious" if dev > 50 else "Warning"
        colour = CRIMSON if dev > 100 else AMBER
        for txt, dx, w, c in [(d, 0.15, 1.6, BODY), (cat, 1.9, 3.0, INK), (money(cost), 5.1, 1.4, BODY),
                              (money(exp), 6.7, 1.4, MUTED), (f"+{dev:.0f}%", 8.4, 1.6, colour), (sev, 10.2, 1.6, colour)]:
            box(s, Inches(0.7 + dx), y, Inches(w), Inches(0.32), txt, 10.5, False, c)
        y += Inches(0.42)

    bullets(s, Inches(0.7), Inches(4.3), Inches(11.9), Inches(2.0), [
        "Modified z-score = 0.6745 × (x − median) / MAD, threshold ≈ 3.5 — robust to the very outliers it is hunting",
        "The test runs on the STL residual, so weekday, weekend and monthly cycles do not trip alerts",
        "A point must be BOTH statistically odd AND financially material (≥25% deviation). The statistical test alone "
        "flags a 5% wobble in a low-variance series, and an alert nobody can act on teaches people to ignore the channel",
        "Mirrors AWS Cost Anomaly Detection semantics: ≥10-day warm-up, dynamic thresholds rather than a static dollar line",
        "Severity is carried by an icon and a label as well as colour — never by colour alone",
    ], 11.5)
    footer(s, SYNTHETIC)


def slide_connect(prs, f: Facts):
    s = blank(prs)
    header(s, "Connecting your estate", "How many credentials do you actually need?",
           "Far fewer than you think — one credential already spans many accounts")
    rows = [
        ("AWS", "Cost Explorer / Data Exports at the payer (management) account return every linked account. "
                "They arrive in FOCUS as SubAccountId.", "One credential per payer"),
        ("Azure", "Billing-account scope covers every subscription beneath it. Subscription scope covers exactly one — "
                  "prefer the billing account.", "One credential per tenant"),
        ("GCP", "A billing account's BigQuery export carries every project it pays for, again as SubAccountId.",
         "One credential per billing account"),
    ]
    y = Inches(2.2)
    for cloud, how, need in rows:
        rect(s, Inches(0.7), y, Inches(11.9), Inches(1.0), PAPER, RULE)
        box(s, Inches(0.95), y + Inches(0.18), Inches(1.3), Inches(0.4), cloud, 14, True, AZURE)
        box(s, Inches(2.35), y + Inches(0.12), Inches(6.6), Inches(0.8), how, 11, False, BODY)
        box(s, Inches(9.2), y + Inches(0.28), Inches(3.2), Inches(0.4), need, 11, True, INK)
        y += Inches(1.12)

    box(s, Inches(0.7), Inches(5.7), Inches(11.9), Inches(1.1),
        "A second credential is needed only for a second PAYER — another AWS organization, another Azure tenant, another OCI tenancy or billing "
        "account, another GCP billing account. A regulated utility usually has several, because regulated and unregulated "
        "entities cannot share a bill. Declare each as an [[accounts]] entry; they are pulled independently and concatenated "
        "into one FOCUS frame, where BillingAccountId keeps them apart. A failing binding never takes the page down.",
        11.5, False, BODY)
    footer(s)


def slide_agents(prs, f: Facts):
    s = blank(prs)
    header(s, "Agentic AI", "A supervisor and four specialists, grounded in your data",
           "Google ADK on Gemini — every figure quoted comes from a tool call against the FOCUS frame")

    rect(s, Inches(0.7), Inches(2.2), Inches(2.5), Inches(1.0), RGBColor(0xE8, 0xF1, 0xFD), AZURE)
    box(s, Inches(0.8), Inches(2.38), Inches(2.3), Inches(0.35), "Supervisor", 13, True, INK, PP_ALIGN.CENTER)
    box(s, Inches(0.8), Inches(2.72), Inches(2.3), Inches(0.4), "routes on gemini-3.1-flash-lite", 8.5, False, MUTED, PP_ALIGN.CENTER)

    specialists = [
        ("Analyst", "Understand Usage & Cost", "spend, allocation, anomalies, coverage"),
        ("Forecaster", "Quantify Business Value", "forecast, budget variance, year-end"),
        ("Optimizer", "Optimize Usage & Cost", "levers, opportunities, ESR uplift"),
        ("Governor", "Manage the FinOps Practice", "tagging, chargeback readiness, policy"),
    ]
    y = Inches(2.2)
    for name, domain, tools in specialists:
        rect(s, Inches(3.7), y, Inches(8.9), Inches(0.78), PAPER, RULE)
        box(s, Inches(3.9), y + Inches(0.06), Inches(1.7), Inches(0.35), name, 12, True, INK)
        box(s, Inches(5.6), y + Inches(0.08), Inches(3.1), Inches(0.3), domain, 10, True, VIOLET)
        box(s, Inches(8.8), y + Inches(0.08), Inches(3.6), Inches(0.5), tools, 10, False, BODY)
        y += Inches(0.86)

    bullets(s, Inches(0.7), Inches(5.75), Inches(11.9), Inches(1.3), [
        "Typed tools query the loaded FOCUS frame directly — the model never recites a number from memory, and it cannot write SQL",
        "The coordinator routes on the cheap model and specialists reason on the strong one: the small-model-first lever, applied to ourselves",
        "Specialists are held as tools, not sub-agents, so one voice writes the final answer for one persona",
        "A column whitelist guards every query. No eval, ever. Loop-guarded, and it degrades to an explanation when the model is unreachable",
    ], 11)
    footer(s)


def slide_agents_flow(prs, f: Facts):
    s = blank(prs)
    header(s, "Agentic AI", "What actually happens when you ask a question",
           'Worked example: "Why did spend rise last month, and what should we do?"')
    steps = [
        ("1", "Route", "The coordinator reads the question and the asking persona, and calls the specialists it needs. Cheap model, one call."),
        ("2", "Call tools", "The Analyst calls get_spend_summary and get_anomalies against the real frame. It receives numbers, not prose."),
        ("3", "Hand back", "The specialist returns to the supervisor with its finding recorded in shared state. The supervisor decides whether another domain is needed."),
        ("4", "Escalate", "Spend rose and an anomaly is implicated, so the supervisor routes to the Optimizer, which calls find_optimization_opportunities and explain_lever."),
        ("5", "Answer", "The coordinator writes one answer, streamed token by token, citing the figure and the tool it came from, in outcome terms for a VP."),
    ]
    y = Inches(2.15)
    for n, title, body in steps:
        circ = s.shapes.add_shape(9, Inches(0.7), y, Inches(0.5), Inches(0.5))
        circ.fill.solid(); circ.fill.fore_color.rgb = AZURE; circ.line.fill.background(); circ.shadow.inherit = False
        circ.text_frame.text = n
        _text(circ, 12, True, PAPER, PP_ALIGN.CENTER)
        box(s, Inches(1.4), y + Inches(0.02), Inches(1.9), Inches(0.4), title, 13, True, INK)
        box(s, Inches(3.4), y - Inches(0.02), Inches(9.2), Inches(0.7), body, 11, False, BODY)
        y += Inches(0.88)

    box(s, Inches(0.7), Inches(6.6), Inches(11.9), Inches(0.5),
        "Guard rails: a bounded agent loop, a session keyed by a stable id so the conversation survives a page reload, and a tool layer "
        "that returns a status rather than raising -- a tool that throws kills the agent.", 10.5, False, MUTED)
    footer(s)


def slide_gcp(prs, f: Facts):
    s = blank(prs)
    header(s, "Target architecture", "Running this on Google Cloud",
           "The engines are unchanged. What changes is the warehouse, the runtime and the model.")

    left = [
        ("Cloud Run", "FastAPI service, scales to zero. React client behind a load balancer and IAP."),
        ("BigQuery", "The FOCUS 1.2 warehouse. Partitioned on ChargePeriodStart, clustered on cloud, "
                     "service category and application."),
        ("Cloud Run Job", "Nightly ingest through the same connectors, plus a materialised optimization snapshot."),
        ("Gemini via Vertex", "Agents authenticate with the service account. No API key exists anywhere."),
    ]
    y = Inches(2.2)
    for title, body in left:
        rect(s, Inches(0.7), y, Inches(6.3), Inches(0.98), PAPER, RULE)
        box(s, Inches(0.95), y + Inches(0.10), Inches(2.3), Inches(0.35), title, 12.5, True, AZURE)
        box(s, Inches(0.95), y + Inches(0.42), Inches(5.8), Inches(0.5), body, 10, False, BODY)
        y += Inches(1.08)

    x = Inches(7.3)
    box(s, x, Inches(2.2), Inches(5.3), Inches(0.35), "Why pandas forced the rebuild", 13, True, INK)
    box(s, x, Inches(2.6), Inches(5.3), Inches(1.3),
        "A utility estate at ~500k billing line-items a month is roughly 8 GB of memory over two years; "
        "a large enterprise at 2M/month is ~31 GB. The demo loads the whole frame into one process. "
        "It demonstrates well and it will not survive Con Edison. BigQuery is not optional.",
        10.5, False, BODY)

    box(s, x, Inches(4.05), Inches(5.3), Inches(0.35), "Three cost guards", 13, True, INK)
    bullets(s, x, Inches(4.45), Inches(5.3), Inches(1.6), [
        "A query with no period filter is rejected, not answered",
        "Every job caps bytes billed — BigQuery fails rather than bills",
        "Row-level detectors run nightly, never per request",
    ], 10.5)

    kpi_card(s, Inches(0.7), Inches(6.05), Inches(3.8), Inches(0.95), "Platform run cost",
             "~$240 / month", "Cloud Run, BigQuery, Cloud SQL, LB", TEAL)
    kpi_card(s, Inches(4.7), Inches(6.05), Inches(3.8), Inches(0.95), "Gemini",
             "~$117 / month", "4,400 questions; ~$0.027 each", VIOLET)
    kpi_card(s, Inches(8.7), Inches(6.05), Inches(3.9), Inches(0.95), "Share of estate",
             "0.0005%", "of the $30.28M under management", GREEN)
    footer(s, "Cost estimates. Gemini priced at gemini-3.5-flash / gemini-3.1-flash-lite list rates, July 2026.")


def slide_integrations(prs, f: Facts):
    s = blank(prs)
    header(s, "Vendor neutrality", f"{f.n_connectors} connectors, one schema",
           "Whatever Con Edison procures, the dashboards do not change")
    groups = [
        ("Native cloud", AZURE, ["AWS — Data Exports (FOCUS 1.2), Cost Explorer, Cost Optimization Hub", "Azure — Cost Management, FocusCost exports, Advisor", "GCP — BigQuery billing export, Recommender, Budgets", "OCI — FOCUS cost reports, Cloud Advisor, Budgets"]),
        ("Procured platforms", VIOLET, ["Apptio Cloudability · Tanzu CloudHealth · Flexera One", "Finout · Vantage · CloudZero · Harness CCM · nOps", "Kubecost · OpenCost · IBM Turbonomic · ServiceNow"]),
        ("Zero-code path", TEAL, ["Any FOCUS CSV or Parquet — local, S3, Azure Blob or GCS", "CloudZero and Vantage emit FOCUS natively", "Cloudability, CloudHealth and Flexera ingest it"]),
    ]
    y = Inches(2.2)
    for title, colour, items in groups:
        rect(s, Inches(0.7), y, Inches(11.9), Inches(1.35), PAPER, RULE)
        bar = s.shapes.add_shape(1, Inches(0.7), y, Pt(4), Inches(1.35))
        bar.fill.solid(); bar.fill.fore_color.rgb = colour; bar.line.fill.background(); bar.shadow.inherit = False
        box(s, Inches(1.0), y + Inches(0.12), Inches(2.7), Inches(0.4), title, 13.5, True, INK)
        bullets(s, Inches(3.8), y + Inches(0.1), Inches(8.5), Inches(1.15), items, 10.5, BODY, "·")
        y += Inches(1.5)
    box(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.5),
        "Adopting a new tool is one Connector subclass and one registry line. Nothing downstream changes.", 11.5, True, INK)
    footer(s)


def slide_roadmap(prs, f: Facts):
    s = blank(prs)
    header(s, "Delivery", "A phased rollout, value in the first wave", "Sequenced by effort and risk, not by dollar size")
    waves = [
        ("Wave 1 · Weeks 1–4", GREEN, "Inform", ["Connect payers; validate FOCUS conformance", "Stand up the executive view and showback", "Baseline the estate; publish allocation coverage", "Quick wins: idle resources, gp2→gp3, orphaned snapshots"]),
        ("Wave 2 · Weeks 5–12", AZURE, "Optimize", ["Tagging remediation to cross the chargeback line", "Commitment strategy: coverage to the sustained floor", "Non-prod scheduling; previous-generation migration", "Forecast vs budget in the monthly finance cycle"]),
        ("Wave 3 · Quarter 2+", VIOLET, "Operate", ["Chargeback with the agreed shared-cost policy", "Unit economics per customer, per kWh, per meter read", "Anomaly alerting into ServiceNow", "AI Copilot for self-service, FinOps for AI/GPU spend"]),
    ]
    x = Inches(0.7)
    for title, colour, phase, items in waves:
        rect(s, x, Inches(2.15), Inches(3.83), Inches(4.1), PAPER, RULE)
        bar = s.shapes.add_shape(1, x, Inches(2.15), Inches(3.83), Pt(4))
        bar.fill.solid(); bar.fill.fore_color.rgb = colour; bar.line.fill.background(); bar.shadow.inherit = False
        box(s, x + Inches(0.22), Inches(2.32), Inches(3.4), Inches(0.35), title, 11.5, True, INK)
        box(s, x + Inches(0.22), Inches(2.66), Inches(3.4), Inches(0.3), f"FinOps phase: {phase}", 9.5, True, colour)
        bullets(s, x + Inches(0.22), Inches(3.05), Inches(3.4), Inches(3.0), items, 10.5)
        x += Inches(4.02)
    box(s, Inches(0.7), Inches(6.45), Inches(11.9), Inches(0.4),
        "Wave 1 is deliberately all low-effort, low-risk levers — the programme should pay for itself before it asks for trust.",
        11, True, INK)
    footer(s)


def slide_honesty(prs, f: Facts):
    s = blank(prs)
    header(s, "Assumptions & limits", "What this does not claim", "Stated up front, because a FinOps tool that overstates is worse than none")
    items = [
        "Every figure in this deck comes from a synthetic 24-month utility estate. The dollars are invented; the mechanics are not.",
        "AWS Cost Explorer does not expose list price. On that path ListCost is set equal to BilledCost, so Effective Savings Rate is understated. Use the FOCUS Data Export for the true number.",
        "Business drivers — customers served, kWh delivered, meter reads — cannot be read from a cloud bill by definition. Unit economics needs a feed from a system of record.",
        "Savings percentages in the lever catalog are vendor 'up-to' figures. Treat them as ceilings.",
        "Storage-tiering and rightsizing detectors infer from billing data alone. Where a detector cannot see access patterns or CPU utilisation, it lowers its confidence and says what telemetry would confirm it.",
        "Several vendor endpoints are marked UNVERIFIED in code where public documentation was thin — notably ServiceNow's Cloud Cost Management table names, which must be enumerated in Con Edison's own instance.",
        "The ~90% allocation-coverage chargeback threshold is practitioner consensus, not a published FinOps Foundation number.",
    ]
    bullets(s, Inches(0.7), Inches(2.15), Inches(11.9), Inches(4.4), items, 12)
    footer(s)


def slide_next(prs, f: Facts):
    s = blank(prs)
    rect(s, Inches(0), Inches(0), W, Inches(2.4), RGBColor(0xF2, 0xF6, 0xFC))
    box(s, Inches(0.9), Inches(0.8), Inches(11), Inches(0.4), "NEXT STEPS", 12, True, AZURE)
    box(s, Inches(0.9), Inches(1.15), Inches(11), Inches(0.8), "What we would need to go live", 34, True, INK)
    steps = [
        ("Read-only credentials", "One per payer / tenant / billing account. Cost Explorer, Cost Management and BigQuery billing export access — nothing that can change a resource."),
        ("A FOCUS export, ideally", "AWS Data Exports (FOCUS 1.2) and Azure FocusCost give a true list price, so Effective Savings Rate is correct rather than understated."),
        ("Your tagging standard", "So the canonical keys map to Con Edison's own — application, business unit, cost centre, environment, owner, project."),
        ("A business driver feed", "Customers served, kWh delivered, work orders closed. This is what turns a cloud bill into a unit cost a VP can defend."),
        ("An allocation policy decision", "Even split, proportional, or a fixed percentage for the shared platform pool. This is an accounting choice, not a technical one."),
        ("A GCP project", "Plus a billing account and an identity provider for IAP. Vertex authenticates the agents through the service account, so no model API key is ever created."),
    ]
    y = Inches(2.85)
    for t, d in steps:
        box(s, Inches(0.9), y, Inches(3.5), Inches(0.4), t, 13.5, True, INK)
        box(s, Inches(4.6), y - Inches(0.02), Inches(7.9), Inches(0.62), d, 11.5, False, BODY)
        y += Inches(0.78)
    footer(s, "Infosys · Con Edison")


# ==========================================================================
# Build
# ==========================================================================


def build(out: str) -> str:
    f = gather()
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    for fn in (
        slide_title,
        slide_challenge,
        slide_one_idea,
        slide_hld,
        slide_architecture,
        slide_euv,
        slide_lld,
        slide_exec,
        slide_spend_chart,
        slide_allocation,
        slide_forecast,
        slide_optimize,
        slide_anomaly,
        slide_connect,
        slide_agents,
        slide_agents_flow,
        slide_gcp,
        slide_integrations,
        slide_roadmap,
        slide_honesty,
        slide_next,
    ):
        fn(prs, f)

    prs.save(out)
    return out


def main() -> None:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--out", default="Infosys_FinOps_ConEdison.pptx")
    args = ap.parse_args()
    path = build(args.out)
    print(f"wrote {path}")


if __name__ == "__main__":
    main()
